"""Build the poster as paper/poster.pptx — A1 landscape (841 x 594 mm).

Follows the Müßig 5-element structure (paper/poster_outline.md):
Einführung -> Methode -> Ergebnis -> Schlussfolgerung (highlighted) -> Referenzen.
Industry partner anonymised. Numbers from the repo reports.

Run: .venv/bin/python paper/_build_poster.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "reports" / "figures"

# palette
ACCENT = RGBColor.from_string("007AFF")
TEXT = RGBColor.from_string("1C1C1E")
SECOND = RGBColor.from_string("8E8E93")
WHITE = RGBColor.from_string("FFFFFF")
LIGHT = RGBColor.from_string("F2F2F7")
HYPO_BG = RGBColor.from_string("E5F0FF")
METHOD = {
    "Z-Score": RGBColor.from_string("1F77B4"),
    "ARIMA": RGBColor.from_string("D62728"),
    "Cluster": RGBColor.from_string("2CA02C"),
    "Autoencoder": RGBColor.from_string("9467BD"),
}

# geometry (cm)
PW, PH = 84.1, 59.4
COL_W = 26.0
C1, C2, C3 = 1.5, 29.05, 56.6
BODY_TOP = 9.6


def _para(tf, runs, *, align=PP_ALIGN.LEFT, space_after=8, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if isinstance(runs, str):
        runs = [(runs, 24, False, TEXT)]
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return p


def add_box(slide, left, top, width, height, *, fill=None, line=None,
            shape=MSO_SHAPE.RECTANGLE, anchor=MSO_ANCHOR.TOP):
    sp = slide.shapes.add_shape(shape, Cm(left), Cm(top), Cm(width), Cm(height))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Cm(0.4))
    return sp, tf


def add_text(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def section_header(slide, left, top, num_label):
    tf = add_text(slide, left, top, COL_W, 1.4)
    _para(tf, [(num_label, 34, True, ACCENT)], first=True)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Cm(PW)
    prs.slide_height = Cm(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ---- header band ----
    add_box(slide, 0, 0, PW, 8.6, fill=LIGHT)
    logo, ltf = add_box(slide, 1.5, 1.2, 9.0, 6.2, fill=WHITE, line=ACCENT,
                        anchor=MSO_ANCHOR.MIDDLE)
    _para(ltf, [("THWS", 40, True, ACCENT)], align=PP_ALIGN.CENTER, first=True)
    _para(ltf, [("Logo-Platzhalter", 14, False, SECOND)], align=PP_ALIGN.CENTER)

    htf = add_text(slide, 11.5, 0.9, 71.0, 7.4, anchor=MSO_ANCHOR.MIDDLE)
    _para(htf, [("KI-gestützte Anomalieerkennung in Smart-Meter-Daten", 54, True, TEXT)],
          align=PP_ALIGN.CENTER, first=True, space_after=6)
    _para(htf, [("Ensemble-Methoden und LLM-Handlungsempfehlungen für das "
                 "Energiemanagement", 30, False, ACCENT)], align=PP_ALIGN.CENTER,
          space_after=8)
    _para(htf, [("Felix Zorn · Jakob Ringel   |   Modul „Nachhaltigkeit und "
                 "Digitalisierung“ · Prof. Dr. M. Müßig · FIW · THWS", 22, False, SECOND)],
          align=PP_ALIGN.CENTER)

    # ================= COLUMN 1 — Einführung + Methode =================
    section_header(slide, C1, BODY_TOP, "1 · Einführung")
    tf = add_text(slide, C1, BODY_TOP + 1.4, COL_W, 6.3)
    _para(tf, "Der Smart-Meter-Rollout (GNDEW 2023) macht 15-min-Lastgänge "
              "gewerblicher Liegenschaften verfügbar. Vermeidbarer Verbrauch bleibt ohne "
              "Analyse unsichtbar.", first=True)
    _para(tf, [("Ziel: ", 24, True, TEXT),
               ("erklärbare, übertragbare Anomalie-Erkennung mit automatisierten "
                "Handlungsempfehlungen.", 24, False, TEXT)])

    _, htf2 = add_box(slide, C1, BODY_TOP + 8.0, COL_W, 5.4, fill=HYPO_BG, line=ACCENT,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE, anchor=MSO_ANCHOR.MIDDLE)
    _para(htf2, [("Hypothese", 22, True, ACCENT)], first=True, space_after=6)
    _para(htf2, [("Ein Ensemble komplementärer Methoden erkennt Lastgang-Anomalien "
                  "zuverlässiger als jede Einzelmethode – und eine lokale LLM-Schicht "
                  "überführt sie in konkrete Maßnahmen.", 23, True, TEXT)])

    section_header(slide, C1, BODY_TOP + 14.2, "2 · Methode")
    tf = add_text(slide, C1, BODY_TOP + 15.6, COL_W, 7.2)
    _para(tf, [("Stichprobe: ", 24, True, TEXT),
               ("22 Baumarktfilialen + 1 Sonderfall, RLM 15-min, 2023–2026.",
                24, False, TEXT)], first=True)
    _para(tf, [("Design: ", 24, True, TEXT),
               ("zeitlicher Train/Test-Split an 2025-01-01.", 24, False, TEXT)])
    _para(tf, [("Vier Detektoren + LLM-Pipeline:", 24, True, TEXT)])

    # method chips row
    chip_y = BODY_TOP + 23.3
    chip_w, gap = 5.95, 0.42
    for i, name in enumerate(METHOD):
        cx = C1 + i * (chip_w + gap)
        _, ctf = add_box(slide, cx, chip_y, chip_w, 1.8, fill=METHOD[name],
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE, anchor=MSO_ANCHOR.MIDDLE)
        _para(ctf, [(name, 18, True, WHITE)], align=PP_ALIGN.CENTER, first=True)

    tf = add_text(slide, C1, chip_y + 2.2, COL_W, 4.0)
    _para(tf, [("RLM → STL / Features → vier Methoden → ", 22, False, TEXT),
               ("Union-Ensemble", 22, True, ACCENT),
               (" → LLM-Empfehlung (Qwen 2.5 7B, lokal über Ollama).", 22, False, TEXT)],
          first=True)

    # ================= COLUMN 2 — Ergebnis =================
    section_header(slide, C2, BODY_TOP, "3 · Ergebnis")
    tf = add_text(slide, C2, BODY_TOP + 1.4, COL_W, 4.4)
    _para(tf, [("Vier Methoden komplementär: ", 24, True, TEXT),
               ("κ ≈ 0 paarweise (max. 0,11) → disjunkte Anomalien, kein Einzelsieger.",
                24, False, TEXT)], first=True)
    _para(tf, [("Precision 100 % ", 24, True, TEXT),
               ("über alle vier Methoden (66 Kandidaten manuell bestätigt).",
                24, False, TEXT)])

    # figure 1 — kappa heatmap (wide)
    f1_top = BODY_TOP + 6.1
    slide.shapes.add_picture(str(FIG / "06_kappa_heatmap.png"), Cm(C2), Cm(f1_top),
                             width=Cm(COL_W))
    cap = add_text(slide, C2, f1_top + 5.1, COL_W, 0.9)
    _para(cap, [("Abb. 1: Paarweise Übereinstimmung (Cohen’s κ).", 15, False, SECOND)],
          first=True)

    # figure 2 — sweep flag rates
    f2_top = f1_top + 6.3
    slide.shapes.add_picture(str(FIG / "06_sweep_flag_rates.png"), Cm(C2 + 4.0),
                             Cm(f2_top), width=Cm(18.0))
    cap = add_text(slide, C2, f2_top + 10.1, COL_W, 0.9)
    _para(cap, [("Abb. 2: Flag-Rate über X — Sweet-Spot bei X = 0,25. "
                 "Autoencoder ~20× schneller als ARIMA.", 15, False, SECOND)], first=True)

    # LLM card
    card_top = f2_top + 11.4
    _, ctf = add_box(slide, C2, card_top, COL_W, 8.2, fill=LIGHT, line=ACCENT,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _para(ctf, [("LLM-Pipeline: 66/66 erfolgreich · 0 Schemafehler · 6,6 s/Anomalie",
                 20, True, ACCENT)], first=True, space_after=8)
    _para(ctf, [("Beispiel — Baumarkt_03, nachts: ", 20, True, TEXT),
                ("72,6 kW statt 8,0 kW erwartet (+807 %), Mehrkosten 31,35 €.",
                 20, False, TEXT)])
    _para(ctf, [("Ursache: HVAC-Nachtabsenkung nicht eingehalten · Schweregrad ", 20,
                 False, TEXT),
                ("hoch", 20, True, METHOD["ARIMA"]),
                (" (Konfidenz 0,85).", 20, False, TEXT)])

    # ================= COLUMN 3 — Schlussfolgerung (HIGHLIGHTED) =================
    _, stf = add_box(slide, C3, BODY_TOP, COL_W, 46.6, fill=ACCENT,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, anchor=MSO_ANCHOR.TOP)
    _para(stf, [("4 · Schlussfolgerung & Ausblick", 34, True, WHITE)], first=True,
          space_after=14)
    _para(stf, [("Ein Ensemble aus vier statistischen Methoden plus LLM-Pipeline "
                 "ermöglicht plausible Anomalie-Erkennung mit konkreten "
                 "Handlungsempfehlungen — erklärbar und ohne pro-Standort-Training.",
                 30, True, WHITE)], space_after=16)
    _para(stf, [("Praxis", 25, True, WHITE),
                (": kostenbewertete, priorisierte Anomalien; Schwellwert-Regler ohne "
                 "Neu-Inferenz.", 25, False, WHITE)], space_after=12)
    _para(stf, [("Nachhaltigkeit", 25, True, WHITE),
                (": lokales LLM, kein Cloud-Abfluss, kein Trainings-Energieverbrauch "
                 "(SDG 7 & 9).", 25, False, WHITE)], space_after=12)
    _para(stf, [("Ausblick", 25, True, WHITE),
                (": Übertragung auf weitere Kategorien (Tankstellen, Bürogebäude) als "
                 "Folgeschritt.", 25, False, WHITE)])

    # ---- footer ----
    add_box(slide, 0, PH - 2.6, PW, 2.6, fill=LIGHT)
    ftf = add_text(slide, 1.5, PH - 2.4, 81.1, 2.2, anchor=MSO_ANCHOR.MIDDLE)
    _para(ftf, [("Referenzen: ", 16, True, TEXT),
                ("Chandola et al. (2009), Anomaly Detection: A Survey · "
                 "Cleveland et al. (1990), STL-Dekomposition.    |    "
                 "Datengrundlage und methodisches Feedback: Industriepartner aus dem "
                 "Energiebereich.", 16, False, SECOND)], first=True)

    out = ROOT / "paper" / "poster.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out.relative_to(ROOT)} (A1 landscape {PW} x {PH} cm)")
